"""Builds UHI-Presentation.pptx.

Thirteen slides in the section order a technical panel expects: problem
understanding, literature review, methodology, architecture, data sources,
algorithm, validation, interface, outcomes, limitations, team, references.

Section titles state what the slide contains rather than making a claim. The
evidence carries the argument; a headline that fills a third of the slide only
crowds it out.

Every headline figure is either read from the pipeline's own output at build
time or listed in FIGURES with the document it came from. The validation table
is computed from the committed grid - see presentation/README.md.

Layout follows the reference deck: content column 0.65 to 12.65 in, eyebrow at
0.38, title at 0.95, standfirst at 1.73, body from 2.20, footer rule at 6.88.
Palette and type follow frontend/style.css so the deck and the dashboard read as
one product.

    python presentation/build_deck.py
"""

from __future__ import annotations

import json
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = Path(__file__).resolve().parent / "assets"
OUT = ROOT / "UHI-Presentation.pptx"

# ---------------------------------------------------------------- palette ---
# frontend/style.css :root tokens. Keep these in step with that file.
# Taken from the layout reference deck: a light, warm editorial system.
BG = RGBColor(0xFF, 0xFF, 0xFF)          # page
SURFACE = RGBColor(0xED, 0xEA, 0xE2)     # card / panel fill, warm grey
SURFACE_3 = RGBColor(0xD3, 0xD0, 0xC8)   # rules and hairlines
FG = RGBColor(0x18, 0x1C, 0x20)          # ink
MUTED = RGBColor(0x5C, 0x63, 0x69)       # secondary text
DIM = RGBColor(0x8A, 0x90, 0x96)         # captions, page numbers
ACCENT = RGBColor(0xDE, 0x5C, 0x2D)      # the one accent - burnt orange
PRIMARY = RGBColor(0x37, 0x64, 0x91)     # secondary, used sparingly
SUCCESS = RGBColor(0x26, 0x77, 0x6F)     # teal, for the forward-looking panels
DANGER = RGBColor(0xB2, 0x3A, 0x2A)      # deeper red, semantic only

# The dashboard's magma ramp, kept as the title rule. It encodes temperature and
# is the one place a data palette appears in the chrome.
RAMP = [
    RGBColor(0x28, 0x10, 0x4A),
    RGBColor(0x5E, 0x17, 0x6C),
    RGBColor(0x94, 0x28, 0x6C),
    RGBColor(0xCC, 0x3E, 0x5B),
    RGBColor(0xF6, 0x6E, 0x5C),
    RGBColor(0xFC, 0xBB, 0x8C),
]

UI = "Aptos"
DATA = "Aptos"
DISPLAY = "Aptos Display"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)

REPO = "https://github.com/ankamteja/urban-heat-island-mitigation-simulation"
LINKS = {
    "repo": REPO,
    "docs": f"{REPO}/tree/main/docs",
    "site": "https://urban-heat-island-mitigation-simula.vercel.app",
    "dashboard": "https://urban-heat-island-mitigation-simula.vercel.app/frontend",
}


# ------------------------------------------------------------------ figures ---
def load_figures() -> dict:
    """Pulls what the pipeline emits; the rest is sourced in FIGURES."""
    release = json.loads((ROOT / "frontend/data/release.json").read_text(encoding="utf-8"))
    constants = json.loads((ROOT / "shared/constants.json").read_text(encoding="utf-8"))

    counts = release["action_counts"]
    treatable = sum(v for k, v in counts.items() if k != "None")

    # The model's own metrics file, so no ML figure on a slide can drift from
    # what the pipeline actually recorded.
    metrics = json.loads(
        (ROOT / "Machine Learning & Prediction/Results/metrics.json").read_text(encoding="utf-8")
    )
    by_key = {(r["split"], r["features"], r["model"]): r for r in metrics["results"]}
    blocked = by_key[("spatial_block", "base", "RandomForest")]
    random_split = by_key[("random_80_20", "base", "RandomForest")]
    imp = sorted(metrics["feature_importances"].items(), key=lambda kv: -kv[1])

    return {
        "cells": release["cell_count"],
        "counts": counts,
        "treatable": treatable,
        "upper_bound_cr": release["total_cost_inr"] / 1e7,
        "release_id": release["release_id"],
        "constants": constants,
        "ml": {
            "target": metrics["target"],
            "blocked": blocked,
            "random": random_split,
            "top_feature": imp[0],
            "n_features": len(metrics["feature_importances"]),
        },
    }


# Figures with no machine-readable source, and where each one comes from.
FIGURES = {
    "mean_lst": "27.0",          # STATUS.md / index.html
    "peak_lst": "33.2",          # index.html
    "mean_ndvi": "0.449",        # Machine Learning & Prediction/README.md
    "r2_blocked": "0.513",       # metrics.json, spatial-block split
    "funded_cr": "9.99",         # Decision-Support/ranking.csv, capped at 10 Cr
    "funded_cells": "249",       # ditto
    # The 10 Cr plan, not the whole programme. Pairing the 249-cell headline
    # with the all-4,157-cells drop was the exact denominator mismatch the
    # dashboard exists to prevent.
    "drop_treated": "1.00",      # dashboard, 249 funded cells at 10 Cr
    "drop_grid": "0.03",         # same plan, averaged over all 8,144 cells
    "drop_all_treated": "0.99",  # if all 4,157 eligible were treated
    "drop_all_grid": "0.51",     # ditto, over the whole grid
    "tests": "133",              # pytest tests/
    "hotspots_before": "530",    # top decile of the observed range
    "hotspots_after": "350",     # same, after the funded plan
    "excluded_green": "3,752",   # STATUS.md, already tree cover
    "excluded_water": "193",     # STATUS.md, 149 water + 44 wetland
}


# ------------------------------------------------------------- primitives ---
HLINK_CLR_NS = "http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor"
HLINK_EXT_URI = "{A12FA001-AC4F-418D-AE19-62706E023703}"


def link_uses_text_colour(run):
    """Stop PowerPoint repainting a link run in theme blue with an underline.

    A run's own <a:solidFill> loses to the theme's hyperlink colour, so setting
    font.color alone has no visible effect. The 2018 hlinkClr extension is the
    supported way to say "use the text colour", which is what keeps the links on
    the deck's palette instead of Office default blue.
    """
    rPr = run._r.get_or_add_rPr()
    rPr.set("u", "none")
    hlink = rPr.find(qn("a:hlinkClick"))
    if hlink is None:
        return
    ext_lst = etree.SubElement(hlink, qn("a:extLst"))
    ext = etree.SubElement(ext_lst, qn("a:ext"))
    ext.set("uri", HLINK_EXT_URI)
    clr = etree.SubElement(ext, f"{{{HLINK_CLR_NS}}}hlinkClr")
    clr.set("val", "tx")


def add_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=None, space_after=0):
    """runs: (body, size, colour, bold, font) or (..., font, link) per chunk.

    A list of chunks is one paragraph; a list of lists is several.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for chunk in para:
            body, size, colour, bold, font = chunk[:5]
            link = chunk[5] if len(chunk) > 5 else None
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = colour
            if link:
                r.hyperlink.address = link
                link_uses_text_colour(r)
    return box


def eyebrow(slide, label):
    text(slide, MARGIN, Inches(0.56), Inches(9), Inches(0.3),
         [(label.upper(), 11, ACCENT, True, DATA)])


def heading(slide, title, sub=None, *, size=40):
    text(slide, MARGIN, Inches(0.98), Inches(11.9), Inches(0.9),
         [(title, size, FG, True, UI)])
    if sub:
        text(slide, MARGIN, Inches(1.86), Inches(10.8), Inches(0.5),
             [(sub, 15, MUTED, False, UI)], line_spacing=1.35)


def footer(slide, n, total, *, link=None, link_label=None):
    if link:
        text(slide, MARGIN, Inches(6.86), Inches(9), Inches(0.3),
             [("↗ ", 9, ACCENT, False, DATA),
              (link_label or link, 9, ACCENT, False, DATA, link)])
    else:
        text(slide, MARGIN, Inches(6.86), Inches(7), Inches(0.3),
             [("URBAN HEAT ISLAND MITIGATION · GUWAHATI", 9, DIM, False, DATA)])
    text(slide, Inches(10.6), Inches(6.86), Inches(2.05), Inches(0.3),
         [(f"{n} / {total}", 9, DIM, False, DATA)], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, *, fill=SURFACE, line=None, line_w=1.0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    if line:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    else:
        box.line.fill.background()
    box.text_frame.text = ""
    return box


def stat(slide, x, y, w, value, label, *, colour=FG, value_pt=34):
    text(slide, x, y, w, Inches(0.62), [(value, value_pt, colour, True, UI)])
    text(slide, x, y + Inches(0.62), w, Inches(0.3),
         [(label.upper(), 9.5, DIM, False, DATA)])


def rule(slide, x, y, w, colour=SURFACE_3):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


# ------------------------------------------------------------- primitives ---
# Geometry follows the layout reference: content column 0.65 → 12.65, eyebrow at
# 0.38, section title at 0.95, standfirst at 1.73, body from 2.20, footer rule
# at 6.88. Titles are 30pt, not 40 — a headline that fills a third of the slide
# leaves no room for the evidence underneath it.
MARGIN = Inches(0.65)
COL_W = Inches(12.0)


def section(slide, n, label, title, standfirst=None):
    """Eyebrow, functional section title, optional one-line standfirst.

    The number comes from the slide index rather than the call site, so removing
    a section cannot leave the rest misnumbered.
    """
    text(slide, MARGIN, Inches(0.38), Inches(7.5), Inches(0.25),
         [(f"{n:02d} — {label.upper()}", 10, ACCENT, True, DATA)])
    text(slide, MARGIN, Inches(0.95), COL_W, Inches(0.65),
         [(title, 30, FG, True, DISPLAY)])
    if standfirst:
        text(slide, Inches(0.67), Inches(1.73), Inches(11.4), Inches(0.55),
             [(standfirst, 13.5, MUTED, False, UI)], line_spacing=1.3)


def page(slide, n, *, link=None, link_label=None):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(6.88), COL_W, Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = SURFACE_3
    ln.line.fill.background()
    ln.shadow.inherit = False

    if link:
        text(slide, MARGIN, Inches(7.02), Inches(9.0), Inches(0.22),
             [("↗ ", 9, ACCENT, False, DATA),
              (link_label or link, 9, ACCENT, False, DATA, link)])
    text(slide, Inches(11.9), Inches(7.02), Inches(0.75), Inches(0.22),
         [(f"{n:02d}", 9, DIM, False, DATA)], align=PP_ALIGN.RIGHT)


def wrapped_lines(body, w_inches, size):
    """Rough line count for a text box. PowerPoint reflows at render time, so
    stacked blocks need an estimate here or a three-line paragraph silently
    overlaps the heading beneath it."""
    # 0.58 em per character, measured against Aptos at these sizes. The earlier
    # 0.50 was optimistic: in a 3-inch column it reported one line where two
    # rendered, and stacked bullets overlapped.
    chars_per_line = max(16, int((w_inches * 96) / (size * 0.58)))
    return max(1, -(-len(body) // chars_per_line))


def para(slide, x, y, w, lead, body, *, size=12.5, lead_colour=None, gap=0.22):
    """A labelled paragraph: bold lead-in, then the substance. Returns the y the
    next block should start at, so a column of these cannot collide."""
    text(slide, x, y, w, Inches(0.24),
         [(lead, size, lead_colour or FG, True, UI)])
    text(slide, x, y + Inches(0.26), w, Inches(0.6),
         [(body, size, MUTED, False, UI)], line_spacing=1.32)
    lines = wrapped_lines(body, w / 914400, size)
    return y + Inches(0.26) + Inches(0.205 * lines) + Inches(gap)


def panel(slide, x, y, w, h, title, *, tint=None):
    card(slide, x, y, w, h, fill=SURFACE)
    tab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.05), h)
    tab.fill.solid()
    tab.fill.fore_color.rgb = tint or ACCENT
    tab.line.fill.background()
    tab.shadow.inherit = False
    text(slide, x + Inches(0.26), y + Inches(0.16), w - Inches(0.52), Inches(0.24),
         [(title.upper(), 10, tint or ACCENT, True, DATA)])


def table(slide, x, y, w, headers, rows, *, widths=None, size=12,
          row_h=0.33, highlight=None):
    n = len(headers)
    widths = widths or [1 / n] * n
    cols = [int(w * f) for f in widths]

    hx = x
    for i, h in enumerate(headers):
        text(slide, hx, y, cols[i], Inches(0.26),
             [(h.upper(), 9, DIM, False, DATA)],
             align=PP_ALIGN.RIGHT if i else PP_ALIGN.LEFT)
        hx += cols[i]

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.26), w, Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = SURFACE_3
    ln.line.fill.background(); ln.shadow.inherit = False

    cy = y + Inches(0.36)
    for r_i, row in enumerate(rows):
        on = highlight is not None and r_i == highlight
        cx = x
        for i, cell in enumerate(row):
            text(slide, cx, cy, cols[i], Inches(0.28),
                 [(str(cell), size, ACCENT if on else (FG if i == 0 else MUTED),
                   on, UI if i == 0 else DATA)],
                 align=PP_ALIGN.RIGHT if i else PP_ALIGN.LEFT)
            cx += cols[i]
        cy += Inches(row_h)
    return cy


# ----------------------------------------------------------------- slides ---
def slide_title(prs, f, n, total):
    s = add_slide(prs)

    seg = W / len(RAMP)
    for i, c in enumerate(RAMP):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(seg * i), 0, int(seg) + 1, Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = c
        bar.line.fill.background(); bar.shadow.inherit = False

    text(s, MARGIN, Inches(0.62), Inches(8.0), Inches(0.25),
         [("SUSTAINABILITY HACKATHON  /  URBAN CLIMATE", 10, MUTED, False, DATA)])

    text(s, MARGIN, Inches(1.05), Inches(11.5), Inches(0.8),
         [("HEATWISE", 44, FG, True, DISPLAY)])
    text(s, MARGIN, Inches(1.78), Inches(11.5), Inches(0.5),
         [("Urban Heat Island Mitigation Simulation", 20, ACCENT, True, DISPLAY)])

    card(s, MARGIN, Inches(2.36), COL_W, Inches(1.06), fill=SURFACE)
    text(s, MARGIN + Inches(0.3), Inches(2.54), Inches(11.3), Inches(0.22),
         [("PROBLEM STATEMENT", 9, ACCENT, True, DATA)])
    text(s, MARGIN + Inches(0.3), Inches(2.8), Inches(11.3), Inches(0.5),
         [("Identify where urban surface heat is most severe in Guwahati, and decide "
           "which mitigation measures to fund first under a fixed municipal budget.",
           14, FG, False, UI)], line_spacing=1.3)

    cols = [
        (f"{f['cells']:,}", "CELLS MAPPED AT 100 M", FG),
        (f"{f['treatable']:,}", "ELIGIBLE AFTER LAND-COVER RULES", FG),
        (f"₹{FIGURES['funded_cr']} Cr", "COMMITTED AT ₹10 CR BUDGET", ACCENT),
        (FIGURES["funded_cells"], "SITES IN THE SHORTLIST", FG),
    ]
    for i, (v, l, c) in enumerate(cols):
        x = MARGIN + Inches(3.0) * i
        text(s, x, Inches(3.72), Inches(2.9), Inches(0.45),
             [(v, 26, c, True, DISPLAY)])
        text(s, x, Inches(4.22), Inches(2.9), Inches(0.22),
             [(l, 8, MUTED, False, DATA)])

    rule(s, MARGIN, Inches(4.68), COL_W)

    text(s, MARGIN, Inches(4.92), Inches(7.0), Inches(0.22),
         [("TEAM & CONTRIBUTIONS", 9, MUTED, False, DATA)])
    y = Inches(5.2)
    for name, role in TEAM:
        text(s, MARGIN, y, Inches(3.3), Inches(0.24), [(name, 12, FG, True, UI)])
        text(s, Inches(4.05), y, Inches(4.4), Inches(0.24), [(role, 12, MUTED, False, UI)])
        y += Inches(0.28)

    text(s, Inches(9.0), Inches(4.92), Inches(3.65), Inches(0.22),
         [("PROJECT", 9, MUTED, False, DATA)])
    text(s, Inches(9.0), Inches(5.2), Inches(3.65), Inches(1.0),
         [[("Live dashboard", 11.5, ACCENT, False, DATA, LINKS["dashboard"])],
          [("Source repository", 11.5, ACCENT, False, DATA, LINKS["repo"])],
          [("Documentation", 11.5, ACCENT, False, DATA, LINKS["docs"])]],
         line_spacing=1.5)

    text(s, Inches(11.9), Inches(7.02), Inches(0.75), Inches(0.22),
         [("01", 9, MUTED, False, DATA)], align=PP_ALIGN.RIGHT)


def slide_problem(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "problem understanding & motivation",
            "Urban heat is a health and energy burden with an uneven footprint",
            "Guwahati, Assam — 8,144 cells of Landsat-derived surface temperature.")

    png = ASSETS / "dash-overview.jpg"
    if png.exists():
        frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.47), Inches(2.21),
                                   Inches(6.21), Inches(3.54))
        frame.fill.background()
        frame.line.color.rgb = SURFACE_3
        frame.line.width = Pt(0.75)
        frame.shadow.inherit = False
        s.shapes.add_picture(str(png), Inches(6.5), Inches(2.24), width=Inches(6.15))

    blocks = [
        ("Cities run hotter than their surroundings.",
         "Concrete, asphalt and roofing absorb and re-radiate more heat than the "
         "vegetation they replaced."),
        ("The cost is health and energy.",
         "Surface heat drives heat stress in exposed populations and raises cooling "
         "demand when the grid is already at peak."),
        ("The footprint is uneven at street scale.",
         "Across the study area temperature spans 21.1–33.2 °C, and neighbouring "
         "100 m cells differ by several degrees."),
        ("Mitigation budgets are finite.",
         "Treating every eligible cell would cost ₹167.5 Cr against an annual "
         "allocation nearer ₹10 Cr, so placement determines the return."),
    ]
    y = Inches(2.34)
    for lead, body in blocks:
        y = para(s, MARGIN, y, Inches(5.5), lead, body, size=12, gap=0.3)

    page(s, n)


def slide_method(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "proposed methodology & technical approach",
            "From thermal imagery to a budget-constrained shortlist",
            "Five stages, each producing an artefact the next stage consumes.")

    steps = [
        ("01", "Screen", "Composite Landsat 8 thermal and optical bands in Earth "
         "Engine; derive LST, NDVI and NDBI onto a 100 m grid of 8,144 cells."),
        ("02", "Gate", "Classify each cell against ESA WorldCover and remove "
         "never-touch classes. 8,144 cells reduce to 4,157 eligible."),
        ("03", "Price", "Assign a measure per eligible cell and cost it as "
         "rate × coverage fraction × cell area, using published municipal rates."),
        ("04", "Rank", "Order by cooling per rupee, breaking ties on surface "
         "temperature so the hottest eligible sites are funded first."),
        ("05", "Cap", "Accumulate cost down the ranking and stop at the budget. "
         "At ₹10 Cr this yields 249 sites."),
    ]
    y = Inches(2.3)
    for num, name, body in steps:
        text(s, MARGIN, y, Inches(0.6), Inches(0.3), [(num, 14, ACCENT, True, DATA)])
        text(s, Inches(1.25), y, Inches(1.6), Inches(0.3), [(name, 13.5, FG, True, UI)])
        text(s, Inches(3.0), y, Inches(9.65), Inches(0.5),
             [(body, 12.5, MUTED, False, UI)], line_spacing=1.3)
        y += Inches(0.7)

    rule(s, MARGIN, Inches(5.9), COL_W)
    text(s, MARGIN, Inches(6.1), Inches(2.4), Inches(0.24),
         [("DISTINGUISHING CHOICES", 10, DIM, False, DATA)])
    text(s, Inches(3.0), Inches(6.06), Inches(9.65), Inches(0.7),
         [("Suitability is a gate applied before ranking, not advice attached after it. "
           "Cost enters the objective rather than a footnote. The machine-learning model "
           "is reported as a screening diagnostic and is deliberately kept out of the "
           "selection path.", 12, MUTED, False, UI)], line_spacing=1.32)

    page(s, n)


def slide_architecture(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "system architecture & processing pipeline",
            "Six layers, and the artefacts that cross between them")

    png = ASSETS / "architecture.png"
    if png.exists():
        # 1.28:1, so it is sized on height and sits in the left column.
        pic_h = Inches(4.62)
        pic_w = Inches(4.62 * 2400 / 1880)
        s.shapes.add_picture(str(png), MARGIN, Inches(2.02),
                             width=int(pic_w), height=int(pic_h))

    x2 = Inches(7.35)
    notes = [
        ("Separation of concerns",
         "Acquisition, suitability and selection are distinct modules with a file "
         "contract between them. Each can be re-run without the others."),
        ("Rules precede ranking",
         "Layer 3 removes ineligible cells before any ordering happens, so an "
         "unsuitable site cannot surface however hot it is."),
        ("Cost enters the objective",
         "The unit cost model runs alongside the gate, making cooling per rupee the "
         "quantity layer 4 optimises."),
        ("The loop is closed",
         "A scheduled monthly Earth Engine refresh re-enters at layer 2, and CI "
         "regenerates every downstream artefact."),
    ]
    y = Inches(2.1)
    for lead, body in notes:
        y = para(s, x2, y, Inches(5.3), lead, body, size=12, gap=0.26)

    page(s, n, link=LINKS["docs"], link_label="docs/01-architecture.md")


def slide_data(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "data sources, preprocessing & screening model",
            "Inputs, the steps applied before analysis, and what the model does")

    text(s, MARGIN, Inches(2.24), Inches(5.9), Inches(0.24),
         [("1.  SOURCE DATASETS", 10, DIM, False, DATA)])
    table(s, MARGIN, Inches(2.56), Inches(5.9),
          ["Dataset", "Res.", "Role"],
          [["Landsat 8 C2 L2", "30 m", "Thermal + optical"],
           ["ESA WorldCover 2021", "10 m", "Land-cover class"],
           ["Municipal rate schedules", "—", "Unit costs"]],
          widths=[.45, .17, .38], row_h=0.34)

    text(s, MARGIN, Inches(4.16), Inches(5.9), Inches(0.24),
         [("2.  PREPROCESSING STEPS", 10, DIM, False, DATA)])
    steps = [
        "Cloud and shadow masking from the QA_PIXEL band",
        "Multi-scene median composite to remove transient artefacts",
        "LST from the thermal band; NDVI and NDBI from surface reflectance",
        "Zonal aggregation onto a 100 m QGIS fishnet — 8,144 cells",
        "Majority land-cover class assigned per cell from WorldCover",
    ]
    y = Inches(4.5)
    for st in steps:
        text(s, MARGIN, y, Inches(0.2), Inches(0.24), [("·", 12.5, ACCENT, True, UI)])
        text(s, MARGIN + Inches(0.2), y, Inches(5.7), Inches(0.3),
             [(st, 12, MUTED, False, UI)], line_spacing=1.3)
        y += Inches(0.36)

    x2 = Inches(7.05)
    text(s, x2, Inches(2.24), Inches(5.6), Inches(0.24),
         [("3.  DERIVED FIELDS PER CELL", 10, DIM, False, DATA)])
    table(s, x2, Inches(2.56), Inches(5.6),
          ["Field", "Range", "Use"],
          [["LST", "21.1 – 33.2 °C", "Heat measure"],
           ["NDVI", "−0.14 – 0.78", "Vegetation"],
           ["NDBI", "−0.68 – 0.18", "Built-up"],
           ["Land cover", "7 classes", "Eligibility"],
           ["Priority", "3 tiers", "Filtering"]],
          widths=[.30, .38, .32], row_h=0.32)

    ml = f["ml"]
    text(s, x2, Inches(4.5), Inches(5.6), Inches(0.24),
         [("4.  SCREENING MODEL", 10, DIM, False, DATA)])
    card(s, x2, Inches(4.8), Inches(5.6), Inches(1.86), fill=SURFACE)
    rows = [
        ["Target", "Land surface temperature"],
        ["Features", f"NDBI, NDVI, vegetation, lat, lon ({ml['n_features']})"],
        ["Model", "RandomForestRegressor, seed 42"],
        ["Validation", "spatial-block split"],
        ["R²", f"{ml['blocked']['r2']:.3f}   ({ml['random']['r2']:.3f} random split)"],
        ["MAE / RMSE", f"{ml['blocked']['mae']:.2f} / {ml['blocked']['rmse']:.2f} °C"],
        ["Strongest driver", f"{ml['top_feature'][0]}, {ml['top_feature'][1]:.2f} importance"],
    ]
    y = Inches(4.98)
    for k, v in rows:
        text(s, x2 + Inches(0.24), y, Inches(1.9), Inches(0.22), [(k, 10.5, MUTED, False, UI)])
        text(s, x2 + Inches(2.2), y, Inches(3.2), Inches(0.22), [(v, 10.5, FG, False, DATA)])
        y += Inches(0.23)

    text(s, x2, Inches(6.72), Inches(5.6), Inches(0.2),
         [("Blocked split quoted: adjacent cells are near-duplicates, a random split leaks.",
           9.5, DIM, False, UI)])

    page(s, n, link=LINKS["docs"], link_label="docs/07-data-contracts.md")


def slide_algorithm(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "selection algorithm & worked example",
            "The cost model, the ordering rule, and one cell traced through it")

    text(s, MARGIN, Inches(2.24), Inches(5.9), Inches(0.24),
         [("COST AND COOLING MODEL", 10, DIM, False, DATA)])
    table(s, MARGIN, Inches(2.56), Inches(5.9),
          ["Measure", "₹/m²", "Cover", "ΔT °C", "₹/cell"],
          [["Cool roof", "300", "15%", "1.0", "4.01 L"],
           ["Tree canopy", "150", "25%", "0.8", "3.75 L"],
           ["Pocket park", "1,150", "10%", "2.0", "11.5 L"]],
          widths=[.30, .16, .16, .16, .22], highlight=0, row_h=0.33)

    text(s, MARGIN, Inches(3.92), Inches(5.9), Inches(0.5),
         [("Cost = rate × coverage fraction × cell area. Rates are published "
           "municipal figures. ΔT values are planning assumptions.",
           11, DIM, False, UI)], line_spacing=1.3)

    text(s, MARGIN, Inches(4.6), Inches(5.9), Inches(0.24),
         [("ORDERING RULE", 10, DIM, False, DATA)])
    card(s, MARGIN, Inches(4.9), Inches(5.9), Inches(1.4), fill=SURFACE)
    text(s, MARGIN + Inches(0.26), Inches(5.1), Inches(5.4), Inches(0.3),
         [("sort by  ΔT / cost  desc", 13, ACCENT, False, DATA)])
    text(s, MARGIN + Inches(0.26), Inches(5.42), Inches(5.4), Inches(0.3),
         [("then LST desc, then grid_id asc", 13, FG, False, DATA)])
    text(s, MARGIN + Inches(0.26), Inches(5.76), Inches(5.4), Inches(0.4),
         [("accumulate cost down the list; stop at the budget ceiling",
           11.5, MUTED, False, UI)], line_spacing=1.25)

    x2 = Inches(7.05)
    text(s, x2, Inches(2.24), Inches(5.6), Inches(0.24),
         [("WORKED EXAMPLE — RANK 1 OF 4,157", 10, DIM, False, DATA)])
    card(s, x2, Inches(2.56), Inches(5.6), Inches(3.74), fill=SURFACE, line=ACCENT, line_w=1.0)
    text(s, x2 + Inches(0.26), Inches(2.78), Inches(5.1), Inches(0.28),
         [("+102070+29080", 14, FG, True, DATA)])
    text(s, x2 + Inches(0.26), Inches(3.08), Inches(5.1), Inches(0.26),
         [("Railway Colony, Maligaon  ·  26.1235, 91.6915", 11, MUTED, False, UI)])
    table(s, x2 + Inches(0.26), Inches(3.5), Inches(5.08),
          ["Attribute", "Value", "Basis"],
          [["Surface temperature", "33.2 °C", "measured"],
           ["Built-up intensity", "0.121", "measured"],
           ["Vegetation index", "0.210", "measured"],
           ["Land cover", "Built-up", "classified"],
           ["Assigned measure", "Cool roof", "rule"],
           ["Estimated cost", "₹4.01 L", "rate table"],
           ["Expected cooling", "−1.00 °C", "assumption"]],
          widths=[.46, .28, .26], size=11.5, row_h=0.3)

    page(s, n)


def feature_group(slide, x, y, w, heading, items, *, size=10.5, tint=None):
    """A titled bullet group that reports the y it finished at, so columns of
    them pack without collision regardless of how each line wraps."""
    text(slide, x, y, w, Inches(0.2),
         [(heading.upper(), 8.5, tint or ACCENT, True, DATA)])
    cy = y + Inches(0.24)
    for it in items:
        text(slide, x, cy, Inches(0.14), Inches(0.2), [("·", size, tint or ACCENT, True, UI)])
        text(slide, x + Inches(0.15), cy, w - Inches(0.15), Inches(0.36),
             [(it, size, MUTED, False, UI)], line_spacing=1.22)
        cy += Inches(0.19) * wrapped_lines(it, (w - Inches(0.15)) / 914400, size)
        cy += Inches(0.05)
    return cy + Inches(0.2)


def shot(slide, path, x, y, w, caption, sub):
    """A framed screenshot with its label above. The frame matters because the
    captures are dark and would otherwise float on the white page."""
    h = int(w * 988 / 1720)
    text(slide, x, y, w, Inches(0.2), [(caption, 10, ACCENT, True, DATA)])
    text(slide, x, y + Inches(0.22), w, Inches(0.2), [(sub, 10.5, MUTED, False, UI)])
    top = y + Inches(0.48)
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - Inches(0.03), top - Inches(0.03),
                                   w + Inches(0.06), h + Inches(0.06))
    frame.fill.background()
    frame.line.color.rgb = SURFACE_3
    frame.line.width = Pt(0.75)
    frame.shadow.inherit = False
    if Path(path).exists():
        slide.shapes.add_picture(str(path), x, top, width=int(w), height=int(h))
    return top + h


def slide_dashboard(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "dashboard & decision interface",
            "One interface, two budget scenarios")

    # 5.2in wide keeps the pair under 3in tall, which is what the feature block
    # below needs; at full width they ran into it.
    w = Inches(5.2)
    shot(s, ASSETS / "dash-scenario-10cr.jpg", Inches(0.75), Inches(1.86), w,
         "SCENARIO A  ·  ₹10 CR", "249 sites funded · hotspot cells 530 → 350")
    shot(s, ASSETS / "dash-scenario-max.jpg", Inches(7.4), Inches(1.86), w,
         "SCENARIO B  ·  ₹167.5 CR", "all 4,157 sites funded · hotspot cells 530 → 113")

    text(s, MARGIN, Inches(5.4), COL_W, Inches(0.22),
         [("The same selection rule under two ceilings — moving the budget slider "
           "re-runs it live, nothing is pre-rendered.", 11.5, MUTED, False, UI)])

    rule(s, MARGIN, Inches(5.66), COL_W)

    cols = [
        ("Map and layers", [
            "Interpolated field or raw 100 m cells",
            "Layer toggles and priority filter",
            "Drag-box area selection, place search",
            "Temperature ramp with scale legend",
        ]),
        ("Planning", [
            "Budget slider, ₹10 Cr to full programme",
            "Per-measure toggles with counts",
            "Greedy fill in pipeline rank order",
            "Unavailable measures show the reason",
        ]),
        ("Cell inspection", [
            "Grid ID, place name, priority tier",
            "Why this cell: LST, NDBI, NDVI ranked",
            "Measure, cooling, cost, plan rank",
            "Street View link and cell extent box",
        ]),
        ("Analysis and output", [
            "Distribution, current against plan",
            "Vegetation against temperature",
            "Cooling efficiency per crore",
            "Methodology drawer, report export",
        ]),
    ]
    x = MARGIN
    for heading, items in cols:
        text(s, x, Inches(5.82), Inches(2.9), Inches(0.2),
             [(heading.upper(), 8.5, ACCENT, True, DATA)])
        cy = Inches(6.04)
        for it in items:
            text(s, x, cy, Inches(0.12), Inches(0.18), [("·", 9.5, ACCENT, True, UI)])
            text(s, x + Inches(0.13), cy, Inches(2.85), Inches(0.2),
                 [(it, 9.5, MUTED, False, UI)])
            cy += Inches(0.2)
        x += Inches(3.1)

    page(s, n, link=LINKS["dashboard"], link_label="live dashboard")


def slide_outcomes(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "expected outcomes & impact",
            "What the shortlist achieves, and how the method transfers")

    panel(s, MARGIN, Inches(2.24), Inches(5.85), Inches(2.5), "Measured outcomes at ₹10 Cr")
    outs = [
        ("Hotspot cells", f"{FIGURES['hotspots_before']} → {FIGURES['hotspots_after']}",
         "a 34% reduction in top-decile cells"),
        ("Sites funded", "249", "of 4,157 eligible"),
        ("Committed cost", f"₹{FIGURES['funded_cr']} Cr", "against a ₹167.5 Cr upper bound"),
        ("Cooling, funded cells", "−1.00 °C", "modelled, per treated cell"),
    ]
    y = Inches(2.68)
    for k, v, note in outs:
        text(s, MARGIN + Inches(0.26), y, Inches(2.5), Inches(0.26),
             [(k, 12, MUTED, False, UI)])
        text(s, MARGIN + Inches(2.8), y, Inches(1.5), Inches(0.26),
             [(v, 12.5, FG, True, DATA)], align=PP_ALIGN.RIGHT)
        text(s, MARGIN + Inches(4.4), y, Inches(1.3), Inches(0.26),
             [(note, 10, DIM, False, UI)])
        y += Inches(0.48)

    panel(s, Inches(7.05), Inches(2.24), Inches(5.6), Inches(2.5),
          "Scalability", tint=SUCCESS)
    scal = [
        "Landsat 8 and ESA WorldCover are global and free — the method ports to any "
        "city with a boundary file.",
        "Unit rates and cooling values are configuration, not code; a new city edits "
        "one JSON file.",
        "Runtime is minutes on a laptop for a city of 8,144 cells.",
    ]
    y = Inches(2.7)
    for sc in scal:
        text(s, Inches(7.31), y, Inches(5.1), Inches(0.5),
             [("· ", 11.5, SUCCESS, True, UI), (sc, 11.5, MUTED, False, UI)],
             line_spacing=1.3)
        y += Inches(0.66)

    text(s, MARGIN, Inches(5.06), COL_W, Inches(0.24),
         [("SUSTAINABILITY RELEVANCE", 10, DIM, False, DATA)])
    text(s, MARGIN, Inches(5.36), COL_W, Inches(1.2),
         [("Targeted surface cooling reduces heat exposure for the population living "
           "in the hottest blocks and lowers peak cooling demand, which is where urban "
           "electricity emissions concentrate. Ordering by cooling per rupee rather "
           "than temperature alone admits more sites under the same ceiling, so the "
           "budget removes 180 top-decile cells instead of 154.",
           12.5, MUTED, False, UI)], line_spacing=1.35)

    page(s, n)


def slide_limits(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "limitations & future work",
            "Known constraints, and the work that would resolve them")

    panel(s, MARGIN, Inches(2.24), Inches(5.85), Inches(3.5), "Current limitations", tint=DANGER)
    lims = [
        ("Cooling values are assumptions.", "0.8 / 1.0 / 2.0 °C per measure come from "
         "a planning catalogue, not local measurement, and they set the ranking."),
        ("One rate is unanchored.", "Cool-roof and park rates cite published schedules; "
         "the tree-canopy rate has no municipal equivalent."),
        ("Surface, not air temperature.", "Landsat measures skin temperature; comfort "
         "and health depend on air temperature."),
        ("A single observation window.", "One thermal composite describes one period, "
         "not a seasonal or diurnal cycle."),
    ]
    y = Inches(2.7)
    for lead, body in lims:
        text(s, MARGIN + Inches(0.26), y, Inches(5.3), Inches(0.24),
             [(lead, 12, FG, True, UI)])
        text(s, MARGIN + Inches(0.26), y + Inches(0.24), Inches(5.3), Inches(0.44),
             [(body, 11.5, MUTED, False, UI)], line_spacing=1.28)
        y += Inches(0.8)

    panel(s, Inches(7.05), Inches(2.24), Inches(5.6), Inches(3.5), "Planned work", tint=SUCCESS)
    future = [
        ("Calibrate cooling locally.", "Compare existing parks and canopy against "
         "matched built-up cells across seasons; replace each constant with a fitted "
         "estimate and a range."),
        ("Propagate uncertainty.", "Re-run selection across a scenario grid of rates "
         "and report which cells remain funded throughout."),
        ("Extend temporal coverage.", "Multi-date composites so the plan reflects a "
         "pattern rather than one scene."),
        ("Weight by population.", "Optimise people-degrees, prioritising schools, "
         "hospitals and dense settlement."),
    ]
    y = Inches(2.7)
    for lead, body in future:
        text(s, Inches(7.31), y, Inches(5.1), Inches(0.24),
             [(lead, 12, FG, True, UI)])
        text(s, Inches(7.31), y + Inches(0.24), Inches(5.1), Inches(0.44),
             [(body, 11.5, MUTED, False, UI)], line_spacing=1.28)
        y += Inches(0.8)

    text(s, MARGIN, Inches(6.0), COL_W, Inches(0.5),
         [("All assumptions and their provenance are recorded in the repository under "
           "docs/08-limitations.md and STATUS.md.", 11.5, DIM, False, UI)])

    page(s, n, link=LINKS["docs"], link_label="docs/08-limitations.md")


def slide_references(prs, f, n, total):
    s = add_slide(prs)
    section(s, n, "appendix: links & references",
            "Project resources and cited work")

    text(s, MARGIN, Inches(2.16), Inches(4.0), Inches(0.22),
         [("PROJECT LINKS", 9, MUTED, False, DATA)])
    # Full width per row: these URLs run to 74 characters and were previously
    # clipped, which makes a link slide worse than useless.
    y = Inches(2.46)
    for label, url in [("Live dashboard", LINKS["dashboard"]),
                       ("Landing page", LINKS["site"]),
                       ("Source repository", LINKS["repo"]),
                       ("Documentation", LINKS["docs"])]:
        text(s, MARGIN, y, Inches(2.3), Inches(0.26), [(label, 12, FG, False, UI)])
        text(s, Inches(3.05), y, Inches(9.6), Inches(0.26),
             [(url, 11.5, ACCENT, False, DATA, url)])
        y += Inches(0.36)

    rule(s, MARGIN, Inches(4.06), COL_W)

    text(s, MARGIN, Inches(4.28), Inches(5.9), Inches(0.22),
         [("DATA SOURCES", 9, MUTED, False, DATA)])
    y = Inches(4.58)
    for src in ["USGS Landsat 8–9 Collection 2 Level-2 Science Products",
                "ESA WorldCover 10 m 2021 v200 (Zanaga et al., 2022)",
                "Telangana Cool Roof Policy 2023–28 — ₹300/m²",
                "Gujarat AMRUT 2.0 garden schedule — ₹1,152–2,250/m²"]:
        text(s, MARGIN, y, Inches(5.9), Inches(0.26), [(src, 11.5, MUTED, False, UI)])
        y += Inches(0.32)

    x2 = Inches(7.05)
    text(s, x2, Inches(4.28), Inches(5.6), Inches(0.22),
         [("METHODS AND TOOLING", 9, MUTED, False, DATA)])
    y = Inches(4.58)
    for ref in REFERENCES:
        text(s, x2, y, Inches(5.6), Inches(0.4), [(ref, 10.5, MUTED, False, UI)],
             line_spacing=1.22)
        y += Inches(0.42)

    page(s, n)


# ------------------------------------------------------------- static copy ---
TEAM_NAME = "Team Heatwise"

# Roles are taken from each contributor's commit history rather than assigned.
TEAM = [
    ("Ankam Charan Teja", "Screening model, decision engine, documentation"),
    ("Yogeshwar S", "Dashboard, presentation, pipeline integration"),
    ("Sriya Sudakshina", "Decision-support module, cost catalogue"),
    ("Nikhil Cheepati", "Remote sensing, Earth Engine export"),
    ("Dandu Uddeep Sri Shourya", "Frontend, literature research"),
]

REFERENCES = [
    "Gorelick, N. et al. Google Earth Engine: planetary-scale geospatial analysis. "
    "Remote Sensing of Environment 202, 2017.",
    "Akbari, H., Pomerantz, M. & Taha, H. Cool surfaces and shade trees to reduce "
    "energy use in urban areas. Solar Energy 70(3), 2001.",
    "Santamouris, M. Cooling the cities — reflective and green roof mitigation "
    "technologies. Solar Energy 103, 2014.",
    "Pedregosa, F. et al. Scikit-learn: machine learning in Python. JMLR 12, 2011.",
    "Roberts, D. R. et al. Cross-validation strategies for data with spatial "
    "structure. Ecography 40, 2017.",
]


# -------------------------------------------------------------------- main ---
BUILDERS = [
    slide_title,
    slide_problem,
    slide_method,
    slide_architecture,
    slide_data,
    slide_algorithm,
    slide_dashboard,
    slide_outcomes,
    slide_limits,
    slide_references,
]


def main():
    figures = load_figures()

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    total = len(BUILDERS)
    for i, build in enumerate(BUILDERS, start=1):
        build(prs, figures, i, total)

    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)} - {total} slides, "
          f"{OUT.stat().st_size / 1e6:.2f} MB")


if __name__ == "__main__":
    main()
